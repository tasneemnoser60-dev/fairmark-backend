"""
rag_indexer.py
==============
Indexes teacher-uploaded material (PDF, DOCX, PPTX) into a FAISS vector store.
One index per course — reused for all students in that course.

Usage:
    from grading.rag_indexer import build_index, load_index

    # Teacher uploads material → build index once
    build_index(material_paths=["slides.pdf", "model_answer.docx"], course_id="multimedia_2025")

    # Grader loads index
    retriever = load_index(course_id="multimedia_2025")
"""

import os
import logging
from pathlib import Path

import torch
from llama_index.core import (
    Settings,
    StorageContext,
    VectorStoreIndex,
    load_index_from_storage,
    Document,
)
from llama_index.embeddings.huggingface import HuggingFaceEmbedding

import fitz  # PyMuPDF for PDF
from docx import Document as DocxDocument
from pptx import Presentation

logger = logging.getLogger(__name__)

# ===== Embedding Model (GPU) =====
Settings.embed_model = HuggingFaceEmbedding(
    model_name="BAAI/bge-small-en-v1.5",
    device=os.getenv("EMBEDDING_DEVICE") or ("cuda" if torch.cuda.is_available() else "cpu")
)

# Where indexes are saved
INDEX_BASE_DIR = Path("saved_index")
INDEX_BASE_DIR.mkdir(exist_ok=True)


# ===== FILE READERS =====

def _read_pdf(path: str) -> str:
    text = ""
    with fitz.open(path) as doc:
        for page in doc:
            text += page.get_text()
    return text


def _read_docx(path: str) -> str:
    doc = DocxDocument(path)
    return "\n".join([p.text for p in doc.paragraphs if p.text.strip()])


def _read_pptx(path: str) -> str:
    prs = Presentation(path)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                text += shape.text + "\n"
    return text


def _read_file(path: str) -> str:
    ext = Path(path).suffix.lower()
    if ext == ".pdf":
        return _read_pdf(path)
    elif ext == ".docx":
        return _read_docx(path)
    elif ext == ".pptx":
        return _read_pptx(path)
    else:
        logger.warning(f"Unsupported file type: {ext}, skipping {path}")
        return ""


# ===== BUILD INDEX =====

def build_index(material_paths: list, course_id: str) -> None:
    """
    Build and save a FAISS index from teacher material.
    Call this once per course when teacher uploads material.

    Args:
        material_paths: list of file paths (PDF, DOCX, PPTX)
        course_id: unique identifier for the course (e.g. "multimedia_2025")
    """
    logger.info(f"Building index for course '{course_id}' from {len(material_paths)} file(s)...")

    documents = []
    for path in material_paths:
        logger.info(f"Reading {path}...")
        text = _read_file(path)
        if text.strip():
            documents.append(Document(
                text=text,
                metadata={"source": Path(path).name}
            ))
        else:
            logger.warning(f"No text extracted from {path}")

    if not documents:
        raise ValueError("No content extracted from uploaded material")

    index = VectorStoreIndex.from_documents(documents)

    # Save index to disk
    save_path = INDEX_BASE_DIR / course_id
    index.storage_context.persist(persist_dir=str(save_path))
    logger.info(f"Index saved to {save_path}")


# ===== LOAD INDEX =====

def load_index(course_id: str, top_k: int = 3):
    """
    Load a saved index and return a retriever.

    Args:
        course_id: same ID used in build_index()
        top_k: number of chunks to retrieve per query

    Returns:
        retriever object
    """
    save_path = INDEX_BASE_DIR / course_id
    if not save_path.exists():
        raise FileNotFoundError(f"No index found for course '{course_id}'. Run build_index() first.")

    logger.info(f"Loading index for course '{course_id}'...")
    storage_context = StorageContext.from_defaults(persist_dir=str(save_path))
    index = load_index_from_storage(storage_context)
    return index.as_retriever(similarity_top_k=top_k)


# ===== RETRIEVE CONTEXT =====

def retrieve_context(retriever, query: str) -> list:
    """
    Retrieve relevant chunks for a query.

    Returns:
        list of node objects with .text and .metadata
    """
    return retriever.retrieve(query)


def retrieve_context_text(retriever, query: str) -> str:
    """
    Retrieve relevant chunks and return as a single string.
    """
    nodes = retrieve_context(retriever, query)
    return "\n\n".join([n.text for n in nodes])
